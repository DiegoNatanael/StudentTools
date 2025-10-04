from fastapi import FastAPI, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel, Field
from typing import List
import docx
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as PptRGBColor
import io
from fastapi.middleware.cors import CORSMiddleware

# --- Pydantic Models ---
class DocxSection(BaseModel):
    header: str
    paragraphs: List[str]

class DocumentContent(BaseModel):
    title: str
    sections: List[DocxSection]
    style: str = "formal"

class PptxSlide(BaseModel):
    title: str
    content: List[str]

class PresentationContent(BaseModel):
    title: str
    slides: List[PptxSlide]
    style: str = "formal"  # Add style field

# --- FastAPI App ---
app = FastAPI(title="Document Generation Backend")

# --- CORS Middleware ---
origins = [
    "https://student-tools-front-end.vercel.app",
    "http://127.0.0.1:5500",
    "http://localhost:5500",
    "null",
]
app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- API Endpoints ---

@app.post("/api/generate/docx")
async def generate_docx(content: DocumentContent):
    try:
        document = docx.Document()
        config_style = content.style
        
        # Set default font
        style = document.styles['Normal']
        
        # === FORMAL STYLE: Corporate/Business Report ===
        if config_style == "formal":
            style.font.name = "Times New Roman"
            style.font.size = Pt(12)
            
            # Add professional title
            title = document.add_heading(content.title, level=0)
            title_run = title.runs[0]
            title_run.font.size = Pt(18)
            title_run.font.color.rgb = RGBColor(0, 51, 102)
            title_run.font.bold = True
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title.paragraph_format.space_after = Pt(24)
            
            for section in content.sections:
                # Section headers with underline
                if section.header:
                    header = document.add_heading(section.header, level=1)
                    header_run = header.runs[0]
                    header_run.font.size = Pt(14)
                    header_run.font.color.rgb = RGBColor(0, 51, 102)
                    header_run.font.bold = True
                    header_run.font.underline = True
                    header.paragraph_format.space_before = Pt(12)
                    header.paragraph_format.space_after = Pt(8)
                
                # Regular paragraphs
                for p_text in section.paragraphs:
                    p = document.add_paragraph(p_text)
                    p.paragraph_format.space_after = Pt(10)
                    p.paragraph_format.line_spacing = 1.15
                
                document.add_paragraph()
            
            # Set margins
            for section in document.sections:
                section.top_margin = Inches(1)
                section.bottom_margin = Inches(1)
                section.left_margin = Inches(1)
                section.right_margin = Inches(1)
        
        # === ACADEMIC STYLE: Scholarly Paper ===
        elif config_style == "academic":
            style.font.name = "Georgia"
            style.font.size = Pt(12)
            
            # Centered title
            title = document.add_heading(content.title, level=0)
            title_run = title.runs[0]
            title_run.font.size = Pt(16)
            title_run.font.color.rgb = RGBColor(0, 0, 0)
            title_run.font.bold = True
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            title.paragraph_format.space_after = Pt(36)
            
            for section in content.sections:
                # Bold section headers
                if section.header:
                    header = document.add_paragraph(section.header)
                    header_run = header.runs[0]
                    header_run.font.size = Pt(14)
                    header_run.font.bold = True
                    header.paragraph_format.space_before = Pt(18)
                    header.paragraph_format.space_after = Pt(12)
                
                # Justified paragraphs with first-line indent
                for p_text in section.paragraphs:
                    p = document.add_paragraph(p_text)
                    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    p.paragraph_format.first_line_indent = Inches(0.5)
                    p.paragraph_format.space_after = Pt(12)
                    p.paragraph_format.line_spacing = 2.0
            
            # Wide margins
            for section in document.sections:
                section.top_margin = Inches(1)
                section.bottom_margin = Inches(1)
                section.left_margin = Inches(1.5)
                section.right_margin = Inches(1.5)
        
        # === MODERN STYLE: Clean & Minimal ===
        elif config_style == "modern":
            style.font.name = "Calibri"
            style.font.size = Pt(11)
            
            # Large title
            title = document.add_heading(content.title, level=0)
            title_run = title.runs[0]
            title_run.font.size = Pt(22)
            title_run.font.color.rgb = RGBColor(41, 128, 185)
            title_run.font.bold = True
            title.alignment = WD_ALIGN_PARAGRAPH.LEFT
            title.paragraph_format.space_after = Pt(24)
            
            for section in content.sections:
                # Colored section headers
                if section.header:
                    header = document.add_paragraph(section.header)
                    header_run = header.runs[0]
                    header_run.font.size = Pt(15)
                    header_run.font.color.rgb = RGBColor(52, 73, 94)
                    header_run.font.bold = True
                    header.paragraph_format.space_before = Pt(16)
                    header.paragraph_format.space_after = Pt(10)
                
                # Clean paragraphs
                for p_text in section.paragraphs:
                    p = document.add_paragraph(p_text)
                    p.paragraph_format.space_after = Pt(10)
                    p.paragraph_format.line_spacing = 1.3
                
                document.add_paragraph()
            
            # Narrow margins
            for section in document.sections:
                section.top_margin = Inches(0.75)
                section.bottom_margin = Inches(0.75)
                section.left_margin = Inches(1)
                section.right_margin = Inches(1)
        
        # === SAVE TO BUFFER ===
        doc_buffer = io.BytesIO()
        document.save(doc_buffer)
        doc_buffer.seek(0)
        filename = f"{content.title.replace(' ', '_')}.docx"
        headers = {'Content-Disposition': f'attachment; filename="{filename}"'}
        return StreamingResponse(doc_buffer, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", headers=headers)
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/generate/pptx")
async def generate_pptx(content: PresentationContent):
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        style = content.style
        
        # === FORMAL STYLE: Black background with white text ===
        if style == "formal":
            # Title Slide
            blank_layout = prs.slide_layouts[6]  # Blank layout
            title_slide = prs.slides.add_slide(blank_layout)
            
            # Black background
            background = title_slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = PptRGBColor(0, 0, 0)
            
            # Title text (centered, white)
            title_box = title_slide.shapes.add_textbox(
                Inches(1), Inches(3), Inches(8), Inches(1.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = content.title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = PptPt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = PptRGBColor(255, 255, 255)
            title_para.alignment = 1  # Center
            
            # Content Slides
            for slide_data in content.slides:
                slide = prs.slides.add_slide(blank_layout)
                
                # Black background
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = PptRGBColor(0, 0, 0)
                
                # Title (white, top)
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
                )
                title_frame = title_box.text_frame
                title_frame.text = slide_data.title
                title_para = title_frame.paragraphs[0]
                title_para.font.size = PptPt(32)
                title_para.font.bold = True
                title_para.font.color.rgb = PptRGBColor(255, 255, 255)
                
                # Content (white bullets)
                content_box = slide.shapes.add_textbox(
                    Inches(1), Inches(2), Inches(8), Inches(4.5)
                )
                text_frame = content_box.text_frame
                text_frame.word_wrap = True
                
                for point in slide_data.content:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = PptPt(20)
                    p.font.color.rgb = PptRGBColor(255, 255, 255)
                    p.space_after = PptPt(12)
        
        # === BUSINESS STYLE: Clean corporate design ===
        elif style == "business":
            # Title Slide
            blank_layout = prs.slide_layouts[6]
            title_slide = prs.slides.add_slide(blank_layout)
            
            # White background with blue accent bar
            background = title_slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = PptRGBColor(255, 255, 255)
            
            # Blue accent bar at top
            accent_bar = title_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(10), Inches(0.5)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = PptRGBColor(0, 102, 204)
            accent_bar.line.color.rgb = PptRGBColor(0, 102, 204)
            
            # Title
            title_box = title_slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(8), Inches(2)
            )
            title_frame = title_box.text_frame
            title_frame.text = content.title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = PptPt(40)
            title_para.font.bold = True
            title_para.font.color.rgb = PptRGBColor(0, 51, 102)
            title_para.alignment = 1
            
            # Content Slides
            for slide_data in content.slides:
                slide = prs.slides.add_slide(blank_layout)
                
                # White background
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = PptRGBColor(255, 255, 255)
                
                # Blue accent bar
                accent_bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0), Inches(10), Inches(0.5)
                )
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = PptRGBColor(0, 102, 204)
                accent_bar.line.color.rgb = PptRGBColor(0, 102, 204)
                
                # Title
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1), Inches(9), Inches(0.7)
                )
                title_frame = title_box.text_frame
                title_frame.text = slide_data.title
                title_para = title_frame.paragraphs[0]
                title_para.font.size = PptPt(28)
                title_para.font.bold = True
                title_para.font.color.rgb = PptRGBColor(0, 51, 102)
                
                # Content
                content_box = slide.shapes.add_textbox(
                    Inches(1), Inches(2.2), Inches(8), Inches(4.5)
                )
                text_frame = content_box.text_frame
                text_frame.word_wrap = True
                
                for point in slide_data.content:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = PptPt(18)
                    p.font.color.rgb = PptRGBColor(51, 51, 51)
                    p.space_after = PptPt(12)
        
        # === CREATIVE STYLE: Colorful with visual elements ===
        elif style == "creative":
            # Title Slide
            blank_layout = prs.slide_layouts[6]
            title_slide = prs.slides.add_slide(blank_layout)
            
            # Gradient-like background (light blue)
            background = title_slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = PptRGBColor(240, 248, 255)
            
            # Decorative circles
            circle1 = title_slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(8), Inches(0.5), Inches(2), Inches(2)
            )
            circle1.fill.solid()
            circle1.fill.fore_color.rgb = PptRGBColor(255, 107, 107)
            circle1.line.fill.background()
            
            circle2 = title_slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.2), Inches(5.5), Inches(1.5), Inches(1.5)
            )
            circle2.fill.solid()
            circle2.fill.fore_color.rgb = PptRGBColor(85, 239, 196)
            circle2.line.fill.background()
            
            # Title
            title_box = title_slide.shapes.add_textbox(
                Inches(1), Inches(2.5), Inches(8), Inches(2)
            )
            title_frame = title_box.text_frame
            title_frame.text = content.title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = PptPt(42)
            title_para.font.bold = True
            title_para.font.color.rgb = PptRGBColor(88, 24, 69)
            title_para.alignment = 1
            
            # Content Slides
            colors = [
                PptRGBColor(255, 107, 107),  # Red
                PptRGBColor(72, 219, 251),   # Cyan
                PptRGBColor(85, 239, 196),   # Green
                PptRGBColor(253, 203, 110),  # Yellow
                PptRGBColor(162, 155, 254),  # Purple
            ]
            
            for idx, slide_data in enumerate(content.slides):
                slide = prs.slides.add_slide(blank_layout)
                
                # Light background
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = PptRGBColor(250, 250, 250)
                
                # Colorful accent shape (left side)
                accent_color = colors[idx % len(colors)]
                accent_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0), Inches(0.8), Inches(0.3), Inches(5)
                )
                accent_shape.fill.solid()
                accent_shape.fill.fore_color.rgb = accent_color
                accent_shape.line.fill.background()
                
                # Small decorative circle
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(8.5), Inches(0.3), Inches(1.2), Inches(1.2)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = accent_color
                circle.line.fill.background()
                
                # Title
                title_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(0.8), Inches(7.5), Inches(0.8)
                )
                title_frame = title_box.text_frame
                title_frame.text = slide_data.title
                title_para = title_frame.paragraphs[0]
                title_para.font.size = PptPt(30)
                title_para.font.bold = True
                title_para.font.color.rgb = PptRGBColor(51, 51, 51)
                
                # Content
                content_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(2), Inches(8.5), Inches(4.8)
                )
                text_frame = content_box.text_frame
                text_frame.word_wrap = True
                
                for point in slide_data.content:
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.font.size = PptPt(18)
                    p.font.color.rgb = PptRGBColor(51, 51, 51)
                    p.space_after = PptPt(14)
        
        # Save presentation
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        headers = {
            'Content-Disposition': f'attachment; filename="{content.title.replace(" ", "_")}.pptx"'
        }
        return StreamingResponse(ppt_buffer, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers=headers)

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"An error occurred: {str(e)}")